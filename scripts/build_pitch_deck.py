from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/Users/voyager/Desktop/spectrum")
OUTPUT_DIR = ROOT / "output"
ASSET_IMAGE = OUTPUT_DIR / "call-workbench-30min-readable.pdf.png"
OUTFILE = OUTPUT_DIR / "Spectrum_Live_Demo_Pitch.pptx"


COLORS = {
    "navy": RGBColor(15, 23, 42),
    "slate": RGBColor(71, 85, 105),
    "muted": RGBColor(100, 116, 139),
    "border": RGBColor(219, 226, 234),
    "bg": RGBColor(248, 250, 252),
    "white": RGBColor(255, 255, 255),
    "blue": RGBColor(59, 130, 246),
    "blue_soft": RGBColor(219, 234, 254),
    "green": RGBColor(34, 197, 94),
    "green_soft": RGBColor(220, 252, 231),
    "amber": RGBColor(245, 158, 11),
    "amber_soft": RGBColor(254, 243, 199),
    "red": RGBColor(239, 68, 68),
    "red_soft": RGBColor(254, 226, 226),
    "purple": RGBColor(139, 92, 246),
    "purple_soft": RGBColor(237, 233, 254),
}


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    font_size=20,
    bold=False,
    color=None,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["navy"]
    return box


def add_card(slide, left, top, width, height, *, fill=None, line=None, radius_shape=None):
    shape = slide.shapes.add_shape(
        radius_shape or MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill or COLORS["white"]
    shape.line.color.rgb = line or COLORS["border"]
    return shape


def add_chip(slide, left, top, width, text, fill, text_color):
    chip = add_card(
        slide,
        left,
        top,
        width,
        Inches(0.42),
        fill=fill,
        line=fill,
        radius_shape=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
    )
    add_textbox(
        slide,
        left + Inches(0.08),
        top + Inches(0.04),
        width - Inches(0.16),
        Inches(0.28),
        text,
        font_size=12,
        bold=True,
        color=text_color,
    )
    return chip


def add_bullet_list(slide, left, top, width, items, *, font_size=20, color=None):
    box = slide.shapes.add_textbox(left, top, width, Inches(3.0))
    frame = box.text_frame
    frame.clear()
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.color.rgb = color or COLORS["navy"]
    return box


def add_title_block(slide, eyebrow, title, subtitle=None):
    add_textbox(
        slide,
        Inches(0.7),
        Inches(0.35),
        Inches(5.0),
        Inches(0.3),
        eyebrow.upper(),
        font_size=11,
        bold=True,
        color=COLORS["muted"],
    )
    add_textbox(
        slide,
        Inches(0.7),
        Inches(0.68),
        Inches(8.8),
        Inches(0.9),
        title,
        font_size=28,
        bold=True,
        color=COLORS["navy"],
    )
    if subtitle:
        add_textbox(
            slide,
            Inches(0.7),
            Inches(1.28),
            Inches(8.6),
            Inches(0.7),
            subtitle,
            font_size=14,
            color=COLORS["slate"],
        )


def cover_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_card(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, fill=COLORS["bg"], line=COLORS["bg"])
    add_chip(slide, Inches(0.7), Inches(0.45), Inches(1.45), "Hackathon demo", COLORS["blue_soft"], COLORS["blue"])
    add_textbox(
        slide,
        Inches(0.7),
        Inches(1.15),
        Inches(5.8),
        Inches(1.2),
        "Spectrum",
        font_size=30,
        bold=True,
    )
    add_textbox(
        slide,
        Inches(0.7),
        Inches(2.0),
        Inches(5.8),
        Inches(1.5),
        "Google Analytics for the human side of AI conversations",
        font_size=24,
        bold=True,
        color=COLORS["navy"],
    )
    add_textbox(
        slide,
        Inches(0.7),
        Inches(3.15),
        Inches(4.8),
        Inches(1.0),
        "We turn a 30-minute AI call into emotional, behavioral, and decision-ready insights.",
        font_size=16,
        color=COLORS["slate"],
    )

    add_card(slide, Inches(7.0), Inches(0.8), Inches(5.7), Inches(5.9), fill=COLORS["white"], line=COLORS["border"])
    slide.shapes.add_picture(str(ASSET_IMAGE), Inches(7.18), Inches(1.0), width=Inches(5.34))

    add_chip(slide, Inches(0.7), Inches(5.4), Inches(1.35), "Confidence", COLORS["green_soft"], COLORS["green"])
    add_chip(slide, Inches(2.15), Inches(5.4), Inches(1.35), "Hesitation", COLORS["amber_soft"], COLORS["amber"])
    add_chip(slide, Inches(3.6), Inches(5.4), Inches(1.4), "Readiness", COLORS["blue_soft"], COLORS["blue"])
    add_chip(slide, Inches(5.1), Inches(5.4), Inches(1.25), "Friction", COLORS["red_soft"], COLORS["red"])


def problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(
        slide,
        "Problem",
        "Transcripts tell you what was said. They miss how the human actually felt.",
        "Analysts reviewing AI calls still have to guess where confidence rose, hesitation spiked, or the conversation went off track.",
    )

    cards = [
        ("Long calls are slow to review", "A 30-minute call still means replaying audio, skimming text, and manually hunting for important moments.", COLORS["blue_soft"]),
        ("Emotion and behavior get lost", "Words alone do not show stress, doubt, engagement, pauses, fillers, or recovery after a tense turn.", COLORS["amber_soft"]),
        ("Teams do not know what to do next", "Without a clear signal, it is hard to decide whether to follow up, coach, escalate, or qualify the call.", COLORS["red_soft"]),
    ]

    top = Inches(2.2)
    for index, (title, body, fill) in enumerate(cards):
        left = Inches(0.7 + index * 4.07)
        add_card(slide, left, top, Inches(3.7), Inches(2.55), fill=fill, line=COLORS["border"])
        add_textbox(slide, left + Inches(0.18), top + Inches(0.18), Inches(3.3), Inches(0.5), title, font_size=18, bold=True)
        add_textbox(slide, left + Inches(0.18), top + Inches(0.82), Inches(3.25), Inches(1.4), body, font_size=13, color=COLORS["slate"])

    add_card(slide, Inches(0.7), Inches(5.35), Inches(12.0), Inches(0.8), fill=COLORS["navy"], line=COLORS["navy"])
    add_textbox(
        slide,
        Inches(0.95),
        Inches(5.58),
        Inches(11.5),
        Inches(0.3),
        "The missing layer is human signal: emotion, hesitation, confidence, engagement, and readiness during the AI conversation.",
        font_size=16,
        bold=True,
        color=COLORS["white"],
    )


def transcript_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(
        slide,
        "Gap",
        "Why transcripts fail for AI + human call review",
        "A transcript is useful context, but it cannot show timing, emotion, or how the AI changed the human response.",
    )

    add_card(slide, Inches(0.8), Inches(2.0), Inches(5.65), Inches(4.5), fill=COLORS["white"], line=COLORS["border"])
    add_textbox(slide, Inches(1.0), Inches(2.2), Inches(2.0), Inches(0.4), "Transcript view", font_size=16, bold=True, color=COLORS["muted"])
    add_textbox(
        slide,
        Inches(1.0),
        Inches(2.8),
        Inches(5.0),
        Inches(2.7),
        '"Yeah, I think this would actually help our onboarding team because right now everything is scattered across calls and notes."',
        font_size=23,
        color=COLORS["navy"],
    )
    add_bullet_list(
        slide,
        Inches(1.0),
        Inches(4.95),
        Inches(4.9),
        [
            "You can read the sentence",
            "You cannot see the confidence rise",
            "You cannot see the pause before answering",
            "You cannot see the AI calming effect",
        ],
        font_size=15,
        color=COLORS["slate"],
    )

    add_card(slide, Inches(6.75), Inches(2.0), Inches(5.75), Inches(4.5), fill=COLORS["blue_soft"], line=COLORS["border"])
    add_textbox(slide, Inches(6.95), Inches(2.2), Inches(3.0), Inches(0.4), "What analysts actually need", font_size=16, bold=True, color=COLORS["blue"])
    add_bullet_list(
        slide,
        Inches(6.95),
        Inches(2.85),
        Inches(5.0),
        [
            "Where did hesitation spike?",
            "When did confidence stabilize?",
            "Did the AI reduce friction or create it?",
            "Is the human moving toward a decision?",
            "What should the team do after this call?",
        ],
        font_size=20,
        color=COLORS["navy"],
    )


def workbench_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(
        slide,
        "Solution",
        "Spectrum is a single-call voice analytics workbench",
        "The product is built around one core question: how did the human sound, behave, and move toward or away from a decision during the AI conversation?",
    )
    add_card(slide, Inches(0.6), Inches(1.95), Inches(8.1), Inches(4.8), fill=COLORS["white"], line=COLORS["border"])
    slide.shapes.add_picture(str(ASSET_IMAGE), Inches(0.8), Inches(2.15), width=Inches(7.7))

    add_card(slide, Inches(8.95), Inches(2.1), Inches(3.6), Inches(4.4), fill=COLORS["white"], line=COLORS["border"])
    add_textbox(slide, Inches(9.15), Inches(2.3), Inches(2.8), Inches(0.4), "Four job paths in one page", font_size=17, bold=True)

    jobs = [
        ("Investigate", "Find the key moments and understand what happened."),
        ("Score", "Measure confidence, hesitation, engagement, and readiness."),
        ("Coach", "Spot struggle patterns and explain where the human got stuck."),
        ("Decide", "Recommend follow-up, escalation, qualification, or risk review."),
    ]
    for idx, (title, body) in enumerate(jobs):
        y = 2.85 + idx * 0.9
        add_chip(slide, Inches(9.15), Inches(y), Inches(1.35), title, COLORS["purple_soft"], COLORS["purple"])
        add_textbox(slide, Inches(10.65), Inches(y + 0.02), Inches(1.65), Inches(0.42), body, font_size=11, color=COLORS["slate"])


def waveform_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(
        slide,
        "Live Demo",
        "The waveform becomes an analytics surface, not just a player",
        "This is the strongest product moment in the demo because it shows emotion, timing, noise, and key moments on the same 30-minute view.",
    )

    add_card(slide, Inches(0.7), Inches(1.95), Inches(7.7), Inches(4.7), fill=COLORS["white"], line=COLORS["border"])
    pic = slide.shapes.add_picture(str(ASSET_IMAGE), Inches(0.9), Inches(2.15), width=Inches(7.3))
    with Image.open(ASSET_IMAGE) as img:
        width, height = img.size
    # Crop toward the top-left so the slide zooms into the waveform, heatband, and playback panel.
    pic.crop_right = 0.28
    pic.crop_bottom = 0.40
    pic.crop_left = 0.0
    pic.crop_top = 0.08

    add_card(slide, Inches(8.7), Inches(2.1), Inches(3.8), Inches(4.45), fill=COLORS["navy"], line=COLORS["navy"])
    add_textbox(slide, Inches(8.95), Inches(2.35), Inches(2.6), Inches(0.35), "What to point at live", font_size=17, bold=True, color=COLORS["white"])
    highlights = [
        "Full 30-minute waveform for fast scanning",
        "Segment color overlay for emotion clusters",
        "Heatband underneath for clean emotional comparison",
        "Noise track for fan hum, typing, pings, or movement",
        "Emoji raw emotion moments like 😰 😌 💪 😬",
    ]
    add_bullet_list(slide, Inches(8.95), Inches(2.85), Inches(3.1), highlights, font_size=15, color=COLORS["white"])


def workflow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(
        slide,
        "Workflow",
        "One call supports four analyst outcomes",
        "Spectrum is not just an emotion viewer. It is a workbench that supports investigation, scoring, coaching, and decisioning from the same session.",
    )
    titles = ["Investigate", "Score", "Coach", "Decide"]
    bodies = [
        "Jump to key moments like hesitation spikes, emotional recovery, or friction peaks.",
        "Generate human-focused metrics like confidence, engagement, hesitation, and readiness.",
        "Explain where the human struggled, what improved mid-call, and where support is needed.",
        "Recommend follow-up, escalation, qualification, or risk flagging based on the call.",
    ]
    fills = [COLORS["blue_soft"], COLORS["green_soft"], COLORS["amber_soft"], COLORS["red_soft"]]
    for idx in range(4):
        left = Inches(0.7 + idx * 3.05)
        add_card(slide, left, Inches(2.2), Inches(2.75), Inches(3.2), fill=fills[idx], line=COLORS["border"])
        add_textbox(slide, left + Inches(0.16), Inches(2.42), Inches(2.1), Inches(0.35), titles[idx], font_size=19, bold=True)
        add_textbox(slide, left + Inches(0.16), Inches(2.95), Inches(2.35), Inches(1.75), bodies[idx], font_size=13, color=COLORS["slate"])

    add_card(slide, Inches(0.7), Inches(5.75), Inches(12.0), Inches(0.6), fill=COLORS["white"], line=COLORS["border"])
    add_textbox(
        slide,
        Inches(0.95),
        Inches(5.92),
        Inches(11.4),
        Inches(0.25),
        "Primary lens: evaluate the human. AI stays visible as context for interruptions, talk balance, timing, calming effect, and friction.",
        font_size=14,
        bold=True,
        color=COLORS["navy"],
    )


def close_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_card(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, fill=COLORS["navy"], line=COLORS["navy"])
    add_chip(slide, Inches(0.8), Inches(0.6), Inches(1.55), "Why it matters", COLORS["blue"], COLORS["white"])
    add_textbox(
        slide,
        Inches(0.8),
        Inches(1.35),
        Inches(8.8),
        Inches(0.9),
        "Spectrum turns AI conversations into something teams can act on.",
        font_size=28,
        bold=True,
        color=COLORS["white"],
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(2.2),
        Inches(6.6),
        Inches(1.2),
        "Sales teams can spot readiness.\nSupport teams can catch frustration.\nAI teams can measure whether the agent builds trust or creates friction.",
        font_size=19,
        color=RGBColor(226, 232, 240),
    )

    add_card(slide, Inches(8.0), Inches(1.45), Inches(4.2), Inches(3.55), fill=COLORS["white"], line=COLORS["white"])
    metric_cards = [
        ("30 min", "reviewed in seconds"),
        ("4 jobs", "investigate, score, coach, decide"),
        ("1 page", "human signal + AI context"),
    ]
    for idx, (value, label) in enumerate(metric_cards):
        y = 1.8 + idx * 1.02
        add_textbox(slide, Inches(8.35), Inches(y), Inches(1.25), Inches(0.35), value, font_size=24, bold=True)
        add_textbox(slide, Inches(9.75), Inches(y + 0.04), Inches(1.9), Inches(0.3), label, font_size=13, color=COLORS["slate"])

    add_textbox(
        slide,
        Inches(0.8),
        Inches(5.55),
        Inches(7.6),
        Inches(0.4),
        "Spectrum: Google Analytics for the human side of AI conversations",
        font_size=18,
        bold=True,
        color=RGBColor(191, 219, 254),
    )


def build():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    cover_slide(prs)
    problem_slide(prs)
    transcript_slide(prs)
    workbench_slide(prs)
    waveform_slide(prs)
    workflow_slide(prs)
    close_slide(prs)

    prs.save(str(OUTFILE))
    print(f"Created {OUTFILE}")


if __name__ == "__main__":
    build()
